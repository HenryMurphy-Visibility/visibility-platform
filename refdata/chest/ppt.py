
import pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from pptx.dml.color import RGBColor as RGB


# Create a new presentation and slide
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Blank slide layout
slide = prs.slides.add_slide(slide_layout)

def add_shape(slide, shape_type, left, top, width, height, text=""):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.text = text
    return shape

def add_image(slide, image_path, left, top, width, height):
    return slide.shapes.add_picture(image_path, left, top, width, height)

def connect_shapes_with_arrow(slide, shape1, shape2):
    connector = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT,
        0, 0, 100, 100  # Dummy values; will be overwritten below
    )
    connector.begin_connect(shape1, 0)
    connector.end_connect(shape2, 2)
    return connector

def stylize_shape(shape, fill_color=None, line_color=None, font_size=None):
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    if font_size:
        shape.text_frame.paragraphs[0].font.size = font_size

def group_shapes(slide, shape_list):
    sld_shapes = slide.shapes._spTree
    group = pptx.shapes.groupshape.GroupShape(sld_shapes, shape_list)
    return group

# Add the "Events" source on the far left
events_shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.CLOUD, Inches(0.5), Inches(3), Inches(2), Inches(1.5), "Events")

# Add the "Primary Store" to the right of Events
primary_store_shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4), Inches(3), Inches(2), Inches(1.5), "Primary Store")

# Connect Events to Primary Store
connect_shapes_with_arrow(slide, events_shape, primary_store_shape)

world_names = ["Need1", "Need2", "Need3", "DB Layer", "Reconciliation"]

# Slightly modify world positions around the Primary Store and rotate them by ~20 degrees to the left
angle_offset = 20  # degrees
radius = 3.5  # distance from the center of the Primary Store
center_x, center_y = Inches(5), Inches(4)  # center of the circle where worlds are placed

import math
world_angles = [0 - angle_offset, 72 - angle_offset, 144 - angle_offset, 216 - angle_offset, 288 - angle_offset]

world_positions = [(center_x + radius * math.cos(math.radians(angle)), center_y + radius * math.sin(math.radians(angle))) for angle in world_angles]


world_shapes = []  # Initialize an empty list to store the shapes for the worlds

for pos, name in zip(world_positions, world_names):
    world_shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, pos[0], pos[1], Inches(2), Inches(1.5), name)
    world_shapes.append(world_shape)  # Store the shape in the list
    # Connect each world to the Primary Store
    connect_shapes_with_arrow(slide, primary_store_shape, world_shape)

# Save the presentation

prs.save("C:/Users/hjmne/PycharmProjects/chest/pptslides/demo_slide.pptx")
